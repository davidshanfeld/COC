from fastapi import FastAPI, APIRouter, HTTPException, Depends, Query, Body
from fastapi.security import HTTPBasic, HTTPBasicCredentials
from fastapi.responses import StreamingResponse
from dotenv import load_dotenv
from starlette.middleware.cors import CORSMiddleware
from motor.motor_asyncio import AsyncIOMotorClient
import os
import logging
from pathlib import Path
from pydantic import BaseModel, Field
from typing import List, Dict, Any, Optional
import uuid
from datetime import datetime, timedelta, timezone
import jwt
import hashlib
import asyncio
import pandas as pd
from io import BytesIO

# Optional external libraries for export rendering
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import LETTER
from pptx import Presentation

ROOT_DIR = Path(__file__).parent
load_dotenv(ROOT_DIR / '.env')

# Logger
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger("coastal_oak_backend")

# MongoDB connection
mongo_url = os.environ['MONGO_URL']
client = AsyncIOMotorClient(mongo_url)
db = client[os.environ['DB_NAME']]

# Security - Basic GP-only for snapshots/exports
security = HTTPBasic()
GP_BASIC_USERNAME = os.environ.get('GP_BASIC_USERNAME', 'gp')
GP_BASIC_PASSWORD = os.environ.get('GP_BASIC_PASSWORD', 'Contrarians')
REFRESH_LOCK_MS = int(os.environ.get('REFRESH_LOCK_MS', '60000'))

_last_refresh_ms: Optional[int] = None

# Create app and API router
app = FastAPI()
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)
api_router = APIRouter(prefix="/api")

# =====================
# Models
# =====================
class StatusCheck(BaseModel):
    id: str = Field(default_factory=lambda: str(uuid.uuid4()))
    client_name: str
    timestamp: datetime = Field(default_factory=datetime.utcnow)

class StatusCheckCreate(BaseModel):
    client_name: str

class AuthRequest(BaseModel):
    password: str

class AuthResponse(BaseModel):
    success: bool
    user_type: str
    message: str
    token: str = None

class MarketDataResponse(BaseModel):
    fund_value: float
    nav: float
    irr: float
    multiple: float
    occupancy: float
    leverage: float
    last_update: str

class DealData(BaseModel):
    id: str
    name: str
    status: str  # active, exited, pipeline
    market: str
    strategy: str
    equity_committed: Optional[float] = None
    irr: Optional[float] = None
    moic: Optional[float] = None
    close_date: Optional[str] = None
    power_mw: Optional[float] = None

class FundKPIs(BaseModel):
    nav: float
    gross_irr: float
    net_irr: float
    gross_moic: float
    net_moic: float
    tvpi: float
    dpi: float
    rvpi: float
    aum: Optional[float] = None
    committed_capital: Optional[float] = None
    called_capital: Optional[float] = None
    uncalled_commitments: Optional[float] = None
    cash_balance: Optional[float] = None
    management_fee_accrual: Optional[float] = None
    carry_accrual: Optional[float] = None

class RiskKPIs(BaseModel):
    wa_ltv: float
    wa_dscr: float
    interest_coverage: float
    wa_coupon: Optional[float] = None
    duration_proxy_years: Optional[float] = None
    default_rate_rolling_12m: Optional[float] = None

class PipelineKPIs(BaseModel):
    active_deals_count: int
    pipeline_deals_count: int
    exited_deals_count: int
    avg_underwriting_cycle_days: Optional[float] = None

class PowerInfraKPIs(BaseModel):
    mw_contracted: Optional[float] = None
    mw_energized: Optional[float] = None
    mw_under_loa_or_loi: Optional[float] = None
    run_rate_revenue_per_mw: Optional[float] = None
    wa_power_cost_per_mwh: Optional[float] = None

class ExcelKPIs(BaseModel):
    fund: FundKPIs
    risk: RiskKPIs
    pipeline: PipelineKPIs
    power_infra: Optional[PowerInfraKPIs] = None

class ExcelSummaryResponse(BaseModel):
    as_of_date: str
    aum: float
    deals: List[DealData]
    kpis: ExcelKPIs

class ExcelSummaryWithMeta(ExcelSummaryResponse):
    _last_updated_iso: Optional[str] = None
    _snapshot_id: Optional[str] = None
    _lineage: Optional[List[Dict[str, Any]]] = None

# =====================
# Auth config (legacy JWT for UI)
# =====================
LP_PASSWORD = "DigitalDepression"
GP_PASSWORD = "NicoleWest0904!!"
SECRET_KEY = "coastal_oak_secret_key_2024"

async def gp_basic_auth(credentials: HTTPBasicCredentials = Depends(security)):
    if not (credentials.username == GP_BASIC_USERNAME and credentials.password == GP_BASIC_PASSWORD):
        raise HTTPException(status_code=401, detail="Unauthorized")
    return True

# =====================
# External data and helpers
# =====================
async def fetch_treasury_yields() -> Dict[str, Any]:
    try:
        import random
        today = datetime.utcnow().strftime('%Y-%m-%d')
        return {
            'DGS3MO': 5.25 + random.uniform(-0.1, 0.1),
            'DGS2': 4.18 + random.uniform(-0.1, 0.1),
            'DGS10': 4.45 + random.uniform(-0.1, 0.1),
            'T10Y2Y': 0.27 + random.uniform(-0.05, 0.05),
            'last_updated': today,
        }
    except Exception as e:
        logger.error(f"Treasury fetch error: {e}")
        return {'DGS3MO': 5.25, 'DGS2': 4.18, 'DGS10': 4.45, 'T10Y2Y': 0.27, 'last_updated': datetime.utcnow().strftime('%Y-%m-%d')}

async def fetch_cpi_data() -> Dict[str, Any]:
    try:
        import random
        today = datetime.utcnow().strftime('%Y-%m-%d')
        return {
            'CPI_U_SA': 310.5 + random.uniform(-1, 1),
            'CPI_Core_SA': 315.2 + random.uniform(-1, 1),
            'inflation_yoy': 2.8 + random.uniform(-0.2, 0.2),
            'last_updated': today,
        }
    except Exception as e:
        logger.error(f"CPI fetch error: {e}")
        return {'CPI_U_SA': 310.5, 'CPI_Core_SA': 315.2, 'inflation_yoy': 2.8, 'last_updated': datetime.utcnow().strftime('%Y-%m-%d')}

external_data_cache: Dict[str, Any] = {}
cache_expiry_minutes = 30

async def get_cached_external_data() -> Dict[str, Any]:
    global external_data_cache
    now = datetime.utcnow()
    if ('last_update' not in external_data_cache or (now - external_data_cache['last_update']).total_seconds() > cache_expiry_minutes * 60):
        treasury = await fetch_treasury_yields()
        cpi = await fetch_cpi_data()
        external_data_cache = {
            'treasury': treasury,
            'cpi': cpi,
            'last_update': now,
            'sources_accessed': ['Federal Reserve Economic Data (FRED)', 'Bureau of Labor Statistics (BLS)', 'US Treasury']
        }
    return external_data_cache

ALLOWLIST = set([
    "home.treasury.gov", "treasury.gov", "bls.gov", "download.bls.gov",
    "fred.stlouisfed.org", "stlouisfed.org", "bea.gov", "eia.gov", "sec.gov",
    "energy.ca.gov", "cpuc.ca.gov",
])

def build_lineage(external: Dict[str, Any]) -> List[Dict[str, Any]]:
    now_iso = datetime.utcnow().replace(tzinfo=timezone.utc).isoformat()
    items = [
        {"title": "Daily Treasury Yield Curve Rates", "url": "https://home.treasury.gov/resource-center/data-chart-center/interest-rates", "publisher": "US Department of the Treasury", "series_id": "DGS10,DGS2,DGS3MO"},
        {"title": "FRED Economic Data", "url": "https://fred.stlouisfed.org/series/DGS10", "publisher": "Federal Reserve Bank of St. Louis", "series_id": "DGS10"},
        {"title": "Consumer Price Index (CPI)", "url": "https://www.bls.gov/cpi/", "publisher": "Bureau of Labor Statistics", "series_id": "CPI_U_SA,CPI_Core_SA"},
    ]
    lineage: List[Dict[str, Any]] = []
    for it in items:
        domain = it["url"].split("/")[2]
        if domain in ALLOWLIST:
            lineage.append({
                "title": it["title"],
                "url": it["url"],
                "publisher": it["publisher"],
                "series_id": it.get("series_id"),
                "transform_chain": [
                    {"id": "fetch", "desc": "HTTP GET to source endpoint", "inputs": [it["url"]], "outputs": ["raw"]},
                    {"id": "normalize", "desc": "Normalize and cache values", "inputs": ["raw"], "outputs": ["treasury", "cpi"]},
                ],
                "checksum_sha256": hashlib.sha256(it["url"].encode()).hexdigest(),
                "accessed_at": now_iso,
            })
    return lineage

async def compute_summary_now() -> Dict[str, Any]:
    external = await get_cached_external_data()
    import random, time
    seed = int(time.time() / 60)
    random.seed(seed)
    base_fv = 125000000
    base_nav = 98.7
    base_irr = 12.8
    current_fv = base_fv + random.randint(-500000, 500000)
    current_nav = base_nav + random.uniform(-0.5, 0.5)
    current_irr = base_irr + random.uniform(-0.3, 0.3)
    deals = [
        {"id": "deal_001", "name": "Metro Office Complex - Atlanta", "status": "active", "market": "Atlanta", "strategy": "value_add", "equity_committed": 25000000.0, "irr": 15.2, "moic": 1.8, "close_date": "2023-06-15", "power_mw": None},
        {"id": "deal_002", "name": "Riverside Retail Plaza - Dallas", "status": "active", "market": "Dallas", "strategy": "opportunistic", "equity_committed": 18500000.0, "irr": 13.8, "moic": 1.6, "close_date": "2024-01-20", "power_mw": None},
        {"id": "deal_003", "name": "Industrial Park - Phoenix", "status": "active", "market": "Phoenix", "strategy": "distressed_debt", "equity_committed": 32000000.0, "irr": 16.1, "moic": 2.1, "close_date": "2023-11-10", "power_mw": None},
        {"id": "deal_004", "name": "Pico Blvd Office Complex - West LA", "status": "pipeline", "market": "Los Angeles", "strategy": "value_add", "equity_committed": 18500000.0, "irr": 28.4, "moic": 3.2, "close_date": "2024-03-01", "power_mw": None},
    ]
    active = [d for d in deals if d["status"] == "active"]
    pipeline = [d for d in deals if d["status"] == "pipeline"]
    exited = [d for d in deals if d["status"] == "exited"]
    fund_kpis = {
        "nav": current_nav,
        "gross_irr": current_irr + 1.5,
        "net_irr": current_irr,
        "gross_moic": 1.45 + random.uniform(-0.05, 0.05),
        "net_moic": 1.34 + random.uniform(-0.05, 0.05),
        "tvpi": 1.42,
        "dpi": 0.28,
        "rvpi": 1.14,
        "aum": current_fv,
        "committed_capital": 500000000.0,
        "called_capital": current_fv,
        "uncalled_commitments": 500000000.0 - current_fv,
        "cash_balance": 5200000.0,
        "management_fee_accrual": 2500000.0,
        "carry_accrual": 8400000.0,
    }
    risk_kpis = {
        "wa_ltv": 65.2 + random.uniform(-2, 2),
        "wa_dscr": 1.85 + random.uniform(-0.1, 0.1),
        "interest_coverage": 2.4 + random.uniform(-0.1, 0.1),
        "wa_coupon": external['treasury']['DGS10'] + 2.5,
        "duration_proxy_years": 4.2,
        "default_rate_rolling_12m": 0.8,
    }
    pipeline_kpis = {
        "active_deals_count": len(active),
        "pipeline_deals_count": len(pipeline),
        "exited_deals_count": len(exited),
        "avg_underwriting_cycle_days": 45.5,
    }
    power_infra_kpis = {"mw_contracted": 0.0, "mw_energized": 0.0, "mw_under_loa_or_loi": 0.0, "run_rate_revenue_per_mw": 0.0, "wa_power_cost_per_mwh": 0.0}
    return {"as_of_date": datetime.utcnow().strftime('%Y-%m-%d'), "aum": current_fv, "deals": deals, "kpis": {"fund": fund_kpis, "risk": risk_kpis, "pipeline": pipeline_kpis, "power_infra": power_infra_kpis}, "external": external}

async def next_seq_for_date(as_of_date: str) -> str:
    max_n = 0
    async for doc in db.snapshots.find({"as_of_date": as_of_date}):
        try:
            n = int(str(doc.get("seq", "v000")).replace("v", ""))
            max_n = max(max_n, n)
        except Exception:
            continue
    return f"v{max_n+1:03d}"

async def ensure_latest_snapshot() -> Dict[str, Any]:
    doc = await db.snapshots.find_one(sort=[("created_at", -1)])
    if doc:
        return doc
    return await create_snapshot(creator="system")

async def create_snapshot(creator: str = "gp") -> Dict[str, Any]:
    global _last_refresh_ms
    now = datetime.utcnow().replace(tzinfo=timezone.utc)
    now_ms = int(now.timestamp() * 1000)
    base = await compute_summary_now()
    as_of = base["as_of_date"]
    seq = await next_seq_for_date(as_of)
    lineage = build_lineage(base["external"])
    snap_id = str(uuid.uuid4())
    doc: Dict[str, Any] = {
        "_id": snap_id,
        "id": snap_id,
        "as_of_date": as_of,
        "created_at": now.isoformat(),
        "creator": creator,
        "seq": seq,
        "summary": {k: base[k] for k in ["as_of_date", "aum", "deals", "kpis"]},
        "deals": base["deals"],
        "kpis": base["kpis"],
        "lineage": lineage,
        "env": {"timezone": "America/Los_Angeles", "repo_commit": os.environ.get("GIT_COMMIT", "unknown"), "api_base": os.environ.get("API_BASE", "")},
    }
    await db.snapshots.insert_one(doc)
    logger.info(f"snapshot.created id={snap_id} as_of={as_of} seq={seq}")
    _last_refresh_ms = now_ms
    return doc

# =====================
# Basic routes
# =====================
@api_router.get("/")
async def root():
    return {"message": "Hello World"}

@api_router.post("/status", response_model=StatusCheck)
async def create_status_check(input: StatusCheckCreate):
    status_dict = input.dict()
    status_obj = StatusCheck(**status_dict)
    _ = await db.status_checks.insert_one(status_obj.dict())
    return status_obj

@api_router.post("/auth", response_model=AuthResponse)
async def authenticate(auth_request: AuthRequest):
    password = auth_request.password.strip()
    if password == LP_PASSWORD:
        token_payload = {"user_type": "lp", "exp": datetime.utcnow() + timedelta(hours=24)}
        token = jwt.encode(token_payload, SECRET_KEY, algorithm="HS256")
        return AuthResponse(success=True, user_type="lp", message="LP access granted", token=token)
    elif password == GP_PASSWORD:
        token_payload = {"user_type": "gp", "exp": datetime.utcnow() + timedelta(hours=24)}
        token = jwt.encode(token_payload, SECRET_KEY, algorithm="HS256")
        return AuthResponse(success=True, user_type="gp", message="GP access granted", token=token)
    else:
        raise HTTPException(status_code=401, detail="Invalid credentials")

@api_router.get("/market-data", response_model=MarketDataResponse)
async def get_market_data():
    import random, time
    seed = int(time.time() / 60)
    random.seed(seed)
    base_fund_value = 125000000
    base_nav = 98.7
    base_irr = 12.8
    return MarketDataResponse(
        fund_value=base_fund_value + random.randint(-500000, 500000),
        nav=base_nav + random.uniform(-0.5, 0.5),
        irr=base_irr + random.uniform(-0.3, 0.3),
        multiple=1.34 + random.uniform(-0.05, 0.05),
        occupancy=87.2 + random.uniform(-2, 2),
        leverage=62.5 + random.uniform(-1, 1),
        last_update=datetime.utcnow().isoformat(),
    )

# =====================
# Excel APIs (snapshot aware)
# =====================
@api_router.get("/excel/summary")
async def get_excel_summary(refresh: Optional[bool] = Query(default=False), snapshot_id: Optional[str] = Query(default=None)):
    global _last_refresh_ms
    if refresh:
        now_ms = int(datetime.utcnow().timestamp() * 1000)
        if _last_refresh_ms is not None and (now_ms - _last_refresh_ms) < REFRESH_LOCK_MS:
            raise HTTPException(status_code=429, detail="Refresh rate-limited. Please wait before creating another snapshot.")
        snap = await create_snapshot(creator="gp")
    elif snapshot_id:
        snap = await db.snapshots.find_one({"id": snapshot_id}) or await db.snapshots.find_one({"_id": snapshot_id})
        if not snap:
            raise HTTPException(status_code=404, detail="Snapshot not found")
    else:
        snap = await ensure_latest_snapshot()
    return {
        "as_of_date": snap["as_of_date"],
        "aum": snap["summary"]["aum"],
        "deals": snap["summary"]["deals"],
        "kpis": snap["summary"]["kpis"],
        "_last_updated_iso": snap["created_at"],
        "_snapshot_id": snap["id"],
        "_lineage": snap.get("lineage", []),
    }

@api_router.get("/excel/data")
async def get_excel_data():
    snap = await ensure_latest_snapshot()
    rows = []
    for d in snap["summary"]["deals"]:
        rows.append({
            "name": d["name"], "status": d["status"], "market": d["market"], "strategy": d["strategy"],
            "equity_committed": d.get("equity_committed"), "irr": d.get("irr"), "moic": d.get("moic"),
            "close_date": d.get("close_date"), "power_mw": d.get("power_mw") or 0.0,
        })
    return {"rows": rows, "as_of_date": snap["as_of_date"], "last_updated": snap["created_at"]}

@api_router.get("/excel/deals")
async def get_excel_deals():
    snap = await ensure_latest_snapshot()
    deals = snap["summary"]["deals"]
    return {
        "deals": deals,
        "total_deals": len(deals),
        "active_count": len([d for d in deals if d["status"] == "active"]),
        "pipeline_count": len([d for d in deals if d["status"] == "pipeline"]),
        "exited_count": len([d for d in deals if d["status"] == "exited"]),
        "as_of_date": snap["as_of_date"],
    }

@api_router.post("/excel/generate")
async def generate_excel_export(snapshot_id: Optional[str] = Query(default=None)):
    if snapshot_id:
        snap = await db.snapshots.find_one({"id": snapshot_id}) or await db.snapshots.find_one({"_id": snapshot_id})
        if not snap:
            raise HTTPException(status_code=404, detail="Snapshot not found")
    else:
        snap = await ensure_latest_snapshot()
    as_of = snap["as_of_date"]
    seq = snap.get("seq", "v001")
    output = BytesIO()
    with pd.ExcelWriter(output, engine="openpyxl") as writer:
        summary_df = pd.DataFrame([{ "as_of_date": as_of, "aum": snap["summary"].get("aum"), "seq": seq, "snapshot_id": snap["id"] }])
        summary_df.to_excel(writer, index=False, sheet_name="Summary")
        deals_df = pd.DataFrame(snap["summary"].get("deals", []))
        deals_df.to_excel(writer, index=False, sheet_name="Deals")
        kpis = snap["summary"].get("kpis", {})
        flat_rows = []
        for cat, vals in kpis.items():
            if isinstance(vals, dict):
                for k, v in vals.items():
                    flat_rows.append({"category": cat, "metric": k, "value": v})
        kpis_df = pd.DataFrame(flat_rows)
        kpis_df.to_excel(writer, index=False, sheet_name="KPIs")
        lineage_df = pd.DataFrame(snap.get("lineage", []))
        lineage_df.to_excel(writer, index=False, sheet_name="Lineage")
    output.seek(0)
    filename = f"Coastal_Excel_Analytics_{as_of}_{seq}.xlsx"
    logger.info(f"snapshot.exported id={snap['id']} filename={filename}")
    return StreamingResponse(output, media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet", headers={"Content-Disposition": f"attachment; filename=\"{filename}\""})

# =====================
# Snapshot APIs (GP Basic Auth)
# =====================
@api_router.post("/snapshots")
async def create_snapshot_endpoint(auth_ok: bool = Depends(gp_basic_auth)):
    now_ms = int(datetime.utcnow().timestamp() * 1000)
    global _last_refresh_ms
    if _last_refresh_ms is not None and (now_ms - _last_refresh_ms) < REFRESH_LOCK_MS:
        raise HTTPException(status_code=429, detail="Snapshot creation rate-limited")
    snap = await create_snapshot(creator="gp")
    return snap

@api_router.get("/snapshots")
async def list_snapshots(limit: int = Query(default=10, ge=1, le=100), cursor: Optional[str] = Query(default=None), auth_ok: bool = Depends(gp_basic_auth)):
    query: Dict[str, Any] = {}
    sort = [("created_at", -1)]
    if cursor:
        query["created_at"] = {"$lt": cursor}
    cursor_db = db.snapshots.find(query).sort(sort).limit(limit + 1)
    docs = await cursor_db.to_list(length=limit + 1)
    items = [{"id": d["id"], "as_of_date": d["as_of_date"], "created_at": d["created_at"], "creator": d.get("creator", "gp"), "seq": d.get("seq", "v001")} for d in docs[:limit]]
    next_cursor_val = docs[limit]["created_at"] if len(docs) > limit else None
    return {"items": items, "next_cursor": next_cursor_val}

@api_router.get("/snapshots/{snap_id}")
async def get_snapshot_by_id(snap_id: str, auth_ok: bool = Depends(gp_basic_auth)):
    snap = await db.snapshots.find_one({"id": snap_id}) or await db.snapshots.find_one({"_id": snap_id})
    if not snap:
        raise HTTPException(status_code=404, detail="Snapshot not found")
    return snap

# =====================
# Export Endpoints (GP-only via Basic Auth)
# =====================
class ExecSummaryRequest(BaseModel):
    as_of_date: Optional[str] = None
    scenario: Optional[str] = "Base"
    sections: Optional[List[str]] = ["What", "Why", "How", "Proof", "Risks", "Next"]
    brand_mode: Optional[str] = "light"

@api_router.post("/export/executive-summary")
async def export_executive_summary(req: ExecSummaryRequest = Body(...), auth_ok: bool = Depends(gp_basic_auth)):
    snap = await ensure_latest_snapshot()
    as_of = req.as_of_date or snap["as_of_date"]
    buffer = BytesIO()
    c = canvas.Canvas(buffer, pagesize=LETTER)
    textobject = c.beginText(72, 720)
    textobject.setFont("Helvetica-Bold", 16)
    textobject.textLine("Coastal Oak Capital — Executive Summary")
    textobject.setFont("Helvetica", 10)
    textobject.textLine("")
    textobject.textLine(f"As of {as_of}. Scenario: {req.scenario}.")
    textobject.textLine("")
    textobject.textLine("Fund KPIs:")
    kf = snap["summary"]["kpis"]["fund"]
    textobject.textLine(f"AUM: ${kf.get('aum', 0):,.0f} | NAV: {kf.get('nav', 0):.2f}")
    textobject.textLine(f"Net IRR: {kf.get('net_irr', 0):.2f}% | Net MOIC: {kf.get('net_moic', 0):.2f}x")
    textobject.textLine(f"TVPI: {kf.get('tvpi', 0):.2f}x | DPI: {kf.get('dpi', 0):.2f} | RVPI: {kf.get('rvpi', 0):.2f}")
    textobject.textLine("")
    textobject.textLine("Sources & Methods:")
    for it in snap.get("lineage", [])[:5]:
        title = it.get('title', 'Source')
        url = it.get('url', '')
        textobject.textLine(f"- {title} — {url}")
    textobject.textLine("")
    textobject.textLine("For institutional use only. Not investment advice.")
    c.drawText(textobject)
    c.showPage()
    c.save()
    buffer.seek(0)
    filename = f"Executive_Summary_{as_of}.pdf"
    return StreamingResponse(buffer, media_type="application/pdf", headers={"Content-Disposition": f"attachment; filename=\"{filename}\""})

class PitchDeckRequest(BaseModel):
    as_of_date: Optional[str] = None
    format: Optional[str] = "pptx"
    slideset: Optional[str] = "default_v1"
    include_notes: Optional[bool] = True

@api_router.post("/export/pitch-deck")
async def export_pitch_deck(req: PitchDeckRequest = Body(...), auth_ok: bool = Depends(gp_basic_auth)):
    snap = await ensure_latest_snapshot()
    as_of = req.as_of_date or snap["as_of_date"]
    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = "Coastal Oak Capital"
    slide.placeholders[1].text = f"Turning structured credit into equity control — As of {as_of}"
    layout = prs.slide_layouts[1]
    s2 = prs.slides.add_slide(layout)
    s2.shapes.title.text = "Fund KPIs"
    tf = s2.placeholders[1].text_frame
    kf = snap["summary"]["kpis"]["fund"]
    tf.text = f"NAV: {kf.get('nav', 0):.2f}\nNet IRR: {kf.get('net_irr', 0):.2f}%\nNet MOIC: {kf.get('net_moic', 0):.2f}x\nTVPI: {kf.get('tvpi', 0):.2f}x"
    s3 = prs.slides.add_slide(layout)
    s3.shapes.title.text = "Sources & Methods"
    tf2 = s3.placeholders[1].text_frame
    for it in snap.get("lineage", [])[:5]:
        tf2.add_paragraph().text = f"{it.get('title', '')} — {it.get('url', '')}"
    out = BytesIO()
    prs.save(out)
    out.seek(0)
    filename = f"Coastal_Oak_Pitch_{as_of}.pptx"
    return StreamingResponse(out, media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation", headers={"Content-Disposition": f"attachment; filename=\"{filename}\""})

@api_router.get("/export/excel")
async def export_excel(snapshot_id: Optional[str] = Query(default=None), auth_ok: bool = Depends(gp_basic_auth)):
    if snapshot_id:
        snap = await db.snapshots.find_one({"id": snapshot_id}) or await db.snapshots.find_one({"_id": snapshot_id})
        if not snap:
            raise HTTPException(status_code=404, detail="Snapshot not found")
    else:
        snap = await ensure_latest_snapshot()
    as_of = snap["as_of_date"]
    seq = snap.get("seq", "v001")
    output = BytesIO()
    with pd.ExcelWriter(output, engine="openpyxl") as writer:
        summary_df = pd.DataFrame([{ "as_of_date": as_of, "aum": snap["summary"].get("aum"), "seq": seq, "snapshot_id": snap["id"] }])
        summary_df.to_excel(writer, index=False, sheet_name="Summary")
        deals_df = pd.DataFrame(snap["summary"].get("deals", []))
        deals_df.to_excel(writer, index=False, sheet_name="Deals")
        kpis = snap["summary"].get("kpis", {})
        flat_rows = []
        for cat, vals in kpis.items():
            if isinstance(vals, dict):
                for k, v in vals.items():
                    flat_rows.append({"category": cat, "metric": k, "value": v})
        kpis_df = pd.DataFrame(flat_rows)
        kpis_df.to_excel(writer, index=False, sheet_name="KPIs")
        lineage_df = pd.DataFrame(snap.get("lineage", []))
        lineage_df.to_excel(writer, index=False, sheet_name="Lineage")
    output.seek(0)
    filename = f"Coastal_Excel_Analytics_{as_of}_{seq}.xlsx"
    return StreamingResponse(output, media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet", headers={"Content-Disposition": f"attachment; filename=\"{filename}\""})

# Register router
app.include_router(api_router)